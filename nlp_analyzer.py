"""
File Nest - NLP Analyzer Module

Uses Gemini 3.1 Flash Lite to suggest folders based on file content.
Falls back to TextBlob analysis when Gemini is unavailable.

Flow:
    1. Upload file → detect type → extract text
    2. Clean and limit text (~1000 words)
    3. Send filename + text to Gemini 3.1 Flash Lite
    4. Return 3 ranked folder suggestions
       └─ Gemini fail → TextBlob

Requirements:
    pip install google-genai pdfplumber python-docx openpyxl python-pptx python-dotenv textblob

Setup (.env):
    GEMINI_API_KEY=your_key_here
"""

import os
import re
import json
from collections import Counter
from typing import List, Dict, Tuple

from google import genai
from google.genai import types


# Gemini setup: load the API key lazily and create clients only when needed.

def _load_api_key() -> str:
    key = os.environ.get("GEMINI_API_KEY")
    if key:
        return key
    try:
        from dotenv import load_dotenv
        load_dotenv()
        key = os.environ.get("GEMINI_API_KEY")
        if key:
            return key
    except ImportError:
        pass
    raise EnvironmentError(
        "\n[File Nest] GEMINI_API_KEY not set.\n"
        "Add it to your .env file:\n"
        "GEMINI_API_KEY=your_key_here\n"
    )


def _get_client() -> genai.Client:
    return genai.Client(api_key=_load_api_key())


def _runtime_env_value(name: str) -> str:
    try:
        from dotenv import dotenv_values
        env_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")
        env_value = dotenv_values(env_path if os.path.exists(env_path) else None).get(name)
        if env_value is not None:
            return str(env_value)
    except Exception:
        pass
    return os.environ.get(name, "")


def _runtime_env_flag(name: str) -> bool:
    return _runtime_env_value(name).strip().lower() in {"1", "true", "yes", "on"}


_GENERIC_FOLDER_NAMES = {"documents", "files", "misc", "general", "uploads", "other"}

_EMOJI_COLORS = {
    "📁": ("#e8b84b", "#fef7dd"), "🗂️": ("#e8b84b", "#fef7dd"),
    "📂": ("#e8b84b", "#fef7dd"), "🗃️": ("#9b87d4", "#ede8f8"),
    "📄": ("#7ec8e3", "#e0f4fb"), "📑": ("#7ecfb3", "#d9f5ec"),
    "📝": ("#f2a65a", "#fff0dc"), "📋": ("#7ec8e3", "#e0f4fb"),
    "📌": ("#e87a7a", "#fde8e8"), "🔖": ("#f5a7c7", "#fce8f3"),

    "📚": ("#e8b84b", "#fef7dd"), "📖": ("#7ec8e3", "#e0f4fb"),
    "🎓": ("#7ecfb3", "#d9f5ec"), "✏️": ("#f5d06b", "#fef7dd"),
    "🖊️": ("#6c757d", "#f0f0f0"), "📐": ("#9b87d4", "#ede8f8"),
    "🧮": ("#e8855a", "#fde8de"), "🧠": ("#f5a7c7", "#fce8f3"),
    "🔬": ("#7ec8e3", "#e0f4fb"), "🧪": ("#52b788", "#d8f3e8"),

    "💼": ("#c49a6c", "#f7eadc"), "🧾": ("#7ecfb3", "#d9f5ec"),
    "📊": ("#e8b84b", "#fef7dd"), "📈": ("#52b788", "#d8f3e8"),
    "📉": ("#e87a7a", "#fde8e8"), "💰": ("#52b788", "#d8f3e8"),
    "🪙": ("#d99a3d", "#fff1d6"), "🏦": ("#7ec8e3", "#e0f4fb"),
    "🛒": ("#f2a65a", "#fff0dc"), "🧺": ("#c49a6c", "#f7eadc"),

    "💻": ("#6c757d", "#f0f0f0"), "🖥️": ("#6c757d", "#f0f0f0"),
    "⌨️": ("#9b87d4", "#ede8f8"), "🖱️": ("#7ec8e3", "#e0f4fb"),
    "🧑‍💻": ("#7ecfb3", "#d9f5ec"), "🛠️": ("#e8855a", "#fde8de"),
    "⚙️": ("#6c757d", "#f0f0f0"), "🔧": ("#7ec8e3", "#e0f4fb"),
    "🧩": ("#9b87d4", "#ede8f8"), "🚀": ("#e8855a", "#fde8de"),

    "🎨": ("#f5a7c7", "#fce8f3"), "🖼️": ("#9b87d4", "#ede8f8"),
    "📸": ("#9b87d4", "#ede8f8"), "🎬": ("#7ec8e3", "#e0f4fb"),
    "🎵": ("#f5a7c7", "#fce8f3"), "🎧": ("#9b87d4", "#ede8f8"),
    "🎤": ("#e87a7a", "#fde8e8"), "🎮": ("#7ecfb3", "#d9f5ec"),
    "🖌️": ("#f4a6b8", "#fdecef"), "✂️": ("#e8855a", "#fde8de"),

    "🏠": ("#7ec8e3", "#e0f4fb"), "🏡": ("#7ecfb3", "#d9f5ec"),
    "🛏️": ("#9b87d4", "#ede8f8"), "🛋️": ("#b09e94", "#f7f4f0"),
    "🍽️": ("#f5d06b", "#fef7dd"), "☕": ("#8b6f5c", "#f3ebe5"),
    "🧁": ("#f5a7c7", "#fce8f3"), "🍵": ("#7ecfb3", "#d9f5ec"),

    "🏥": ("#e87a7a", "#fde8e8"), "💊": ("#7ec8e3", "#e0f4fb"),
    "🩺": ("#52b788", "#d8f3e8"), "💪": ("#e8855a", "#fde8de"),
    "🧘": ("#9b87d4", "#ede8f8"), "🍎": ("#d94f70", "#fde6ee"),
    "🥗": ("#52b788", "#d8f3e8"), "💧": ("#7ec8e3", "#e0f4fb"),

    "✈️": ("#7ec8e3", "#e0f4fb"), "🧳": ("#b09e94", "#f7f4f0"),
    "🗺️": ("#7ecfb3", "#d9f5ec"), "🏖️": ("#f5d06b", "#fef7dd"),
    "🏕️": ("#52b788", "#d8f3e8"), "🚗": ("#e8855a", "#fde8de"),
    "🚲": ("#7ecfb3", "#d9f5ec"), "🎟️": ("#f5a7c7", "#fce8f3"),

    "⚖️": ("#9b87d4", "#ede8f8"), "🛡️": ("#52b788", "#d8f3e8"),
    "🔐": ("#e8b84b", "#fef7dd"), "🔑": ("#f5d06b", "#fef7dd"),
    "📜": ("#c49a6c", "#f7eadc"), "🤝": ("#7ecfb3", "#d9f5ec"),
    "📣": ("#e8855a", "#fde8de"), "📥": ("#b09e94", "#f7f4f0"),

    "🌸": ("#f5a7c7", "#fce8f3"), "🌷": ("#f4a6b8", "#fdecef"),
    "🌻": ("#f5d06b", "#fef7dd"), "🌿": ("#52b788", "#d8f3e8"),
    "✨": ("#e8b84b", "#fef7dd"), "⭐": ("#f5d06b", "#fef7dd"),
    "☁️": ("#7ec8e3", "#e0f4fb"), "🌙": ("#9b87d4", "#ede8f8"),
}

def _color_for_emoji(emoji: str) -> str:
    return _EMOJI_COLORS.get(emoji, ("#7ec8e3", "#e0f4fb"))[0]

def _bg_for_emoji(emoji: str) -> str:
    return _EMOJI_COLORS.get(emoji, ("#7ec8e3", "#e0f4fb"))[1]

# Text extraction: normalize supported document formats into plain text.
def detect_and_extract(path: str, filename: str) -> Tuple[str, str]:
    """Detect file type and extract text. Returns (file_type, raw_text)."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    extractors = {
        "pdf": _extract_pdf, "docx": _extract_docx, "doc": _extract_docx,
        "xlsx": _extract_xlsx, "xls": _extract_xlsx, "csv": _extract_csv,
        "pptx": _extract_pptx, "ppt": _extract_pptx, "txt": _extract_txt,
    }

    if ext not in extractors:
        return "non-text", ""

    raw_text = normalize_text(extractors[ext](path))
    return ("empty", "") if not raw_text.strip() else ("text", raw_text)


def _extract_pdf(path: str) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            texts = [p.extract_text() for p in pdf.pages if p.extract_text()]
        if texts:
            return " ".join(texts)
    except Exception:
        pass
    try:
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            return " ".join(p.extract_text() for p in reader.pages if p.extract_text())
    except Exception:
        return ""


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.extend(c.text.strip() for c in row.cells if c.text.strip())
        return " ".join(parts)
    except Exception:
        return ""


def _extract_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        cells = [
            str(cell).strip()
            for ws in wb.worksheets
            for row in ws.iter_rows(values_only=True)
            for cell in row
            if cell is not None and str(cell).strip()
        ]
        return " ".join(cells)
    except Exception:
        return ""


def _extract_csv(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = [
            str(shape.text).strip()
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and str(shape.text).strip()
        ]
        return " ".join(texts)
    except Exception:
        return ""


def _extract_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


# Text utilities: clean extraction output before keywording or prompting.
def normalize_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", " ")
    return re.sub(r"\s+", " ", text).strip()


def limit_text(raw_text: str, max_words: int = 300) -> str:
    words = raw_text.split()
    return raw_text.strip() if len(words) <= max_words else " ".join(words[:max_words])


def clean_filename_text(filename: str) -> str:
    base = os.path.splitext(filename)[0]
    base = base.lower().replace("_", " ").replace("-", " ").replace(".", " ")
    return re.sub(r"\s+", " ", re.sub(r"\d+", " ", base)).strip()


def extract_keywords(text: str, max_keywords: int = 8) -> List[str]:
    if not text:
        return []
    words = re.findall(r"[A-Za-z][A-Za-z0-9\-]{2,}", text.lower())
    freq = Counter(words)
    ranked = sorted(freq.items(), key=lambda x: (-x[1], -len(x[0]), x[0]))
    return [w for w, _ in ranked[:max_keywords]]


# JSON repair: salvage valid suggestions from a partial Gemini response.
def _repair_json(raw: str) -> str:
    """Salvage valid JSON from a truncated or fence-wrapped Gemini response."""
    raw = re.sub(r"^```(?:json)?", "", raw.strip(), flags=re.IGNORECASE).strip()
    raw = re.sub(r"```$", "", raw).strip()

    try:
        json.loads(raw)
        return raw
    except json.JSONDecodeError:
        pass

    complete_objects, depth, start = [], 0, None
    for i, ch in enumerate(raw):
        if ch == "{":
            if depth == 0:
                start = i
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0 and start is not None:
                try:
                    complete_objects.append(json.loads(raw[start:i + 1]))
                except json.JSONDecodeError:
                    pass
                start = None

    if complete_objects:
        print(f"WARNING: JSON repaired: salvaged {len(complete_objects)} object(s).")
        return json.dumps(complete_objects)

    raise ValueError("Could not repair malformed JSON from Gemini.")


# Gemini suggestions: ask for ranked folder matches, then validate the result.
def ask_gemini_suggestions(
    filename: str,
    text: str,
    existing_folders: List[str],
) -> List[Dict]:
    """Query Gemini for folder suggestions, with TextBlob fallback on failure."""
    filename_context = clean_filename_text(filename)
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "unknown"
    existing_folders_str = "\n".join(f"- {f}" for f in existing_folders) or "None yet."

    response_schema = {
        "type": "object",
        "required": ["file_summary", "suggestions"],
        "properties": {
            "file_summary": {"type": "string"},
            "suggestions": {
                "type": "array",
                "minItems": 3, "maxItems": 3,
                "items": {
                    "type": "object",
                    "required": ["folder_name", "emoji", "is_new", "confidence", "reason"],
                    "properties": {
                        "folder_name": {"type": "string"},
                        "emoji":       {"type": "string"},
                        "is_new":      {"type": "boolean"},
                        "confidence":  {"type": "integer", "minimum": 0, "maximum": 100},
                        "reason":      {"type": "string"},
                    },
                },
            },
        },
    }

    summary_context = limit_text(text, 220)
    prompt = f"""Summarize this file and suggest exactly 3 folders. Return ONLY a JSON object, no other text.

File: {filename} (.{ext})
Name: {filename_context}
Content: {summary_context if summary_context.strip() else "(none - use filename/type)"}
Existing folders: {existing_folders_str}

Rules:
- file_summary: one plain paragraph, factual, max 120 words.
- Suggest exactly 3 folders.
- Folder names must be specific, Title Case, 2-4 words, no special characters.
- Do not use generic names: Documents, Files, Misc, General, Uploads, Other, School, Work.
- Option 1: most specific.
- Option 2: broader but still specific.
- Option 3: alternative angle.
- Reuse existing folder only if clearly relevant. Copy exact name, emoji and set is_new false.
- If no close match exists, create a new unique folder and set is_new true.
- Avoid repeating folder names unless the topic is truly the same.
- confidence: 85-100 strong, 65-84 good, 40-64 reasonable.
- reason: 8 words or fewer.
- emoji: one related emoji.

{{"file_summary":"...","suggestions":[{{"folder_name":"...","emoji":"...","is_new":true,"confidence":85,"reason":"..."}}]}}"""


    try:
        if _runtime_env_flag("DISABLE_GEMINI"):
            raise RuntimeError("Gemini temporarily disabled by DISABLE_GEMINI")

        client = _get_client()
        response = client.models.generate_content(
            model="gemini-3.1-flash-lite",
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.2,
                top_p=0.85,
                max_output_tokens=2048,
                response_mime_type="application/json",
                response_schema=response_schema,
            ),
        )

        raw = ""
        try:
            raw = response.text
        except Exception:
            try:
                raw = response.candidates[0].content.parts[0].text
            except Exception:
                pass

        raw = (raw or "").strip()
        if _runtime_env_flag("DEBUG_NLP"):
            print(f"\n===== GEMINI RAW RESPONSE =====\n{raw}\n================================\n")

        if not raw:
            raise ValueError("Empty response from Gemini")

        parsed = json.loads(_repair_json(raw))
        if isinstance(parsed, dict):
            file_summary = _clean_file_summary(parsed.get("file_summary"), filename, text)
            items = parsed.get("suggestions", [])
        elif isinstance(parsed, list):
            file_summary = _fallback_file_summary(filename, text)
            items = parsed
        else:
            raise ValueError("Gemini did not return a JSON object")
        if not isinstance(items, list):
            raise ValueError("Gemini did not return folder suggestions")

        suggestions = []
        for item in items[:5]:
            folder_name = str(item.get("folder_name", "")).strip()
            if not folder_name or folder_name.lower() in _GENERIC_FOLDER_NAMES:
                continue

            try:
                confidence = min(100, max(0, int(item.get("confidence", 60))))
            except (TypeError, ValueError):
                confidence = 60

            matched = next((f for f in existing_folders if f.lower() == folder_name.lower()), None)
            if matched:
                folder_name = matched
                is_new = False
                confidence = max(confidence, 80)
            else:
                is_new = bool(item.get("is_new", True))

            emoji = str(item.get("emoji", "📁")).strip() or "📁"
            suggestions.append({
                "folder":     folder_name,
                "emoji":      emoji,
                "color":      _color_for_emoji(emoji),
                "bg":         _bg_for_emoji(emoji),
                "confidence": float(confidence),
                "is_new":     is_new,
                "ai_source":  "gemini",
                "reason":     str(item.get("reason", "")).strip() or "Suggested based on file content.",
            })

        suggestions.sort(key=lambda x: x["confidence"], reverse=True)

        if len(suggestions) < 3:
            for fb in _textblob_fallback_suggestions(filename, text, existing_folders):
                if len(suggestions) >= 3:
                    break
                if not any(s["folder"].lower() == fb["folder"].lower() for s in suggestions):
                    suggestions.append(fb)

        return {"ranked": suggestions[:3], "summary": file_summary}

    except Exception as e:
        print(f"GEMINI ERROR: {e}\nSwitching to TextBlob fallback...")
        suggestions = _textblob_fallback_suggestions(filename, text, existing_folders)

        for s in suggestions:
            s["ai_source"] = "textblob"

        return {"ranked": suggestions, "summary": _fallback_file_summary(filename, text)}


# TextBlob fallback: produce deterministic suggestions when Gemini is disabled.

def _clean_file_summary(summary: str, filename: str, text: str) -> str:
    cleaned = normalize_text(str(summary or ""))
    words = cleaned.split()
    if len(words) > 200:
        cleaned = " ".join(words[:200])
    return cleaned or _fallback_file_summary(filename, text)


def _fallback_file_summary(filename: str, text: str) -> str:
    cleaned = normalize_text(text)
    if cleaned:
        sentences = re.split(r"(?<=[.!?])\s+", cleaned)
        summary_words = []
        for sentence in sentences:
            sentence_words = sentence.split()
            if not sentence_words:
                continue
            if len(summary_words) + len(sentence_words) > 200:
                remaining = 200 - len(summary_words)
                summary_words.extend(sentence_words[:remaining])
                break
            summary_words.extend(sentence_words)
            if len(summary_words) >= 45:
                break
        if summary_words:
            return " ".join(summary_words[:200])

    name_context = clean_filename_text(filename)
    if name_context:
        return f"This file appears to be related to {name_context}."
    return "No readable summary could be generated for this file."

_TB_STOPWORDS = {
    "the", "and", "for", "are", "was", "with", "this", "that", "from",
    "have", "has", "had", "not", "but", "all", "can", "its", "will",
    "one", "been", "also", "into", "than", "more", "your", "our",
    "their", "about", "which", "when", "there", "they", "some", "use",
    "used", "using", "each", "file", "data", "due", "per", "via", "etc",
    "pdf", "doc", "txt", "csv", "total", "amount", "date", "name", "type",
    "item", "list", "line", "page", "note", "ref", "see", "may", "new",
    "get", "set", "put", "let", "any", "two", "top", "end", "sub", "pre",
    "out", "off", "add", "run", "how", "way", "day", "ago", "old", "key",
    "yes", "now", "own", "such", "over", "both", "well", "long", "just",
    "only", "even", "back", "then", "here", "those", "these", "what",
    "after", "before", "where", "while", "every", "other", "during",
    "through", "between", "example", "october", "november", "december",
    "january", "february", "march", "april", "june", "july", "august",
    "september", "rendered", "given", "first", "second", "third", "fourth",
    "fifth", "last", "next", "prev", "previous", "info", "information",
    "number", "version", "draft", "copy", "folder", "document", "upload",
    "misc", "general", "other", "various", "attachment", "untitled",
}

_TB_SKIP_SINGULAR_ENDINGS = (
    "sis", "tics", "ics", "ness", "ment", "tion", "sion",
    "ism", "ogy", "ship", "hood", "ware", "work",
)

_TB_DOMAIN_SCORES = {
    "invoice": 3, "report": 3, "contract": 3, "proposal": 3,
    "research": 3, "thesis": 3, "project": 3, "budget": 3,
    "medical": 3, "health": 3, "legal": 3, "finance": 3,
    "marketing": 3, "design": 3, "development": 2, "software": 2,
    "engineering": 2, "education": 2, "training": 2, "policy": 2,
    "security": 2, "compliance": 2, "analysis": 2, "strategy": 2,
    "management": 2, "operations": 2, "customer": 2, "product": 2,
    "science": 2, "academic": 2, "business": 2, "technical": 2,
}


def _tb_safe_singularize(word: str) -> str:
    if word.endswith(_TB_SKIP_SINGULAR_ENDINGS):
        return word
    try:
        from textblob import Word
        singular = str(Word(word).singularize())
        if not re.search(r"[aeiou]", singular) or len(singular) < max(3, len(word) * 0.5):
            return word
        return singular
    except Exception:
        return word


def _textblob_fallback_suggestions(
    filename: str,
    text: str,
    existing_folders: List[str],
) -> List[Dict]:
    print("TextBlob fallback: analysing filename and content...")
    try:
        combined = f"{clean_filename_text(filename)} {text}".strip()
        raw_words = re.findall(r"[A-Za-z][A-Za-z0-9\-]{2,}", combined.lower())

        processed = [
            s for w in raw_words
            if w not in _TB_STOPWORDS and len(w) >= 4
            for s in [_tb_safe_singularize(w)]
            if s not in _TB_STOPWORDS and len(s) >= 4
        ]

        if not processed:
            print("   TextBlob found no usable terms; using default TextBlob suggestions.")
            return _textblob_default_suggestions(existing_folders)

        freq = Counter(processed)
        top_words = [
            w for w, _ in sorted(
                freq.items(),
                key=lambda kv: freq[kv[0]] * (len(kv[0]) + _TB_DOMAIN_SCORES.get(kv[0], 0)),
                reverse=True,
            )[:12]
        ]
        print(f"   Top TextBlob terms: {top_words[:6]}")

        candidates = []
        if len(top_words) >= 2:
            candidates.append((f"{top_words[0].title()} {top_words[1].title()}", 72.0,
                                f"Top terms from file: '{top_words[0]}', '{top_words[1]}'."))
        if len(top_words) >= 4:
            candidates.append((f"{top_words[2].title()} {top_words[3].title()}", 62.0,
                                f"Secondary terms: '{top_words[2]}', '{top_words[3]}'."))
        elif len(top_words) >= 3:
            candidates.append((f"{top_words[1].title()} {top_words[2].title()}", 58.0,
                                "Recurring content terms in file."))
        if top_words:
            candidates.append((f"{top_words[0].title()} Files", 50.0,
                                f"Grouped by dominant term: '{top_words[0]}'."))

        suggestions, seen = [], set()
        for folder_name, confidence, reason in candidates:
            key = folder_name.lower()
            if key in seen:
                continue
            seen.add(key)
            matched = next((f for f in existing_folders if f.lower() == key), None)
            if matched:
                folder_name, is_new, confidence = matched, False, max(confidence, 78.0)
            else:
                is_new = True
            emoji = _emoji_for_phrase(folder_name)
            suggestions.append({
                "folder": folder_name, "emoji": emoji,
                "color": _color_for_emoji(emoji), "bg": _bg_for_emoji(emoji),
                "confidence": confidence, "is_new": is_new, "reason": reason,
                "ai_source": "textblob",
            })

        suggestions.sort(key=lambda x: x["confidence"], reverse=True)

        for default in _textblob_default_suggestions(existing_folders):
            if len(suggestions) >= 3:
                break
            if not any(s["folder"].lower() == default["folder"].lower() for s in suggestions):
                suggestions.append(default)

        print(f"   TextBlob suggestions: {[s['folder'] for s in suggestions[:3]]}")
        return suggestions[:3]

    except Exception as e:
        print(f"WARNING: TextBlob fallback error: {e}")
        return _textblob_default_suggestions(existing_folders)


def _textblob_default_suggestions(existing_folders: List[str]) -> List[Dict]:
    defaults = [
        ("Inbox", "📥", 30.0, "TextBlob could not identify strong terms."),
        ("General Files", "📁", 25.0, "Default TextBlob suggestion."),
        ("Review Later", "📋", 20.0, "Needs manual review before sorting."),
    ]
    suggestions = []
    for folder_name, emoji, confidence, reason in defaults:
        matched = next((f for f in existing_folders if f.lower() == folder_name.lower()), None)
        final_name = matched or folder_name
        suggestions.append({
            "folder": final_name, "emoji": emoji,
            "color": _color_for_emoji(emoji), "bg": _bg_for_emoji(emoji),
            "confidence": confidence, "is_new": matched is None,
            "reason": reason, "ai_source": "textblob",
        })
    return suggestions


def _emoji_for_phrase(name: str) -> str:
    name_l = name.lower()

    checks = [
        # Default / File organization
        (["inbox", "received", "incoming"], "📥"),
        (["folder", "file", "document", "directory"], "📁"),
        (["collection", "organized", "category", "group"], "🗂️"),
        (["archive", "backup", "compressed", "zip"], "🗂️"),
        (["form", "paper", "sheet"], "📄"),
        (["checklist", "list", "task list"], "📋"),
        (["note", "notes", "memo", "write"], "🗒️"),
        (["important", "pin", "priority", "reminder"], "📌"),
        (["bookmark", "reference", "saved", "favorite"], "🔖"),

        # School / Academic
        (["school", "assignment", "student", "lesson", "quiz", "course", "module", "activity"], "📚"),
        (["research", "study", "paper", "thesis", "journal", "literature"], "📖"),
        (["certificate", "diploma", "award", "graduation"], "🎓"),
        (["math", "calculation", "algebra", "equation"], "📐"),
        (["science", "experiment", "laboratory", "chemistry", "biology"], "🧪"),
        (["analysis", "brainstorm", "idea", "thinking"], "🧠"),

        # Technology / Coding
        (["system", "software", "api", "database", "python", "code", "program", "app", "website"], "💻"),
        (["computer", "desktop", "monitor", "windows", "device"], "🖥️"),
        (["settings", "configuration", "setup"], "⚙️"),
        (["bug", "fix", "repair", "debug", "issue"], "🔧"),
        (["tool", "maintenance", "build"], "🛠️"),
        (["feature", "component", "module", "integration"], "🧩"),
        (["launch", "deployment", "release", "startup"], "🚀"),

        # Business / Finance
        (["business", "office", "work", "company", "employee"], "💼"),
        (["receipt", "billing", "transaction", "record"], "🧾"),
        (["report", "dashboard", "analytics", "statistics", "metrics"], "📊"),
        (["growth", "increase", "progress", "performance"], "📈"),
        (["decrease", "loss", "decline", "drop"], "📉"),
        (["finance", "budget", "payment", "invoice", "tax", "salary", "money"], "💰"),
        (["coin", "income", "savings", "expense"], "💵"),
        (["bank", "account", "loan", "deposit"], "🏦"),
        (["sales", "shopping", "purchase", "order"], "🛍️"),

        # Media / Creative
        (["design", "ui", "ux", "prototype", "wireframe", "layout", "mockup", "art"], "🎨"),
        (["image", "photo", "picture", "screenshot"], "📸"),
        (["video", "movie", "recording", "clip"], "🎬"),
        (["music", "audio", "song", "sound"], "🎵"),
        (["podcast", "headset", "listening"], "🎧"),
        (["presentation", "speech", "voice", "record"], "🎤"),
        (["game", "gaming", "play"], "🎮"),
        (["edit", "draw", "paint"], "🖌️"),

        # Home / Personal
        (["home", "house", "family", "personal"], "🏠"),
        (["room", "bedroom", "sleep", "rest"], "🛏️"),
        (["food", "meal", "dining", "restaurant"], "🍱"),
        (["coffee", "drink", "cafe"], "☕"),
        (["tea", "wellness", "relax"], "🍵"),

        # Health
        (["health", "medical", "hospital", "patient", "doctor"], "🏥"),
        (["medicine", "drug", "tablet", "prescription"], "💊"),
        (["checkup", "clinic", "nurse", "treatment"], "🩺"),
        (["fitness", "exercise", "workout", "gym"], "💪"),
        (["meditation", "mental", "calm", "peace"], "🧘"),
        (["diet", "nutrition", "healthy", "vegetable"], "🥗"),
        (["water", "hydration"], "💧"),

        # Travel
        (["travel", "flight", "trip", "airport"], "✈️"),
        (["luggage", "bag", "packing"], "🧳"),
        (["map", "location", "direction", "place"], "🗺️"),
        (["beach", "vacation", "summer"], "🏖️"),
        (["camp", "camping", "hiking", "nature"], "🏕️"),
        (["car", "drive", "vehicle"], "🚗"),
        (["bike", "cycling", "bicycle"], "🚲"),
        (["ticket", "event", "pass"], "🎟️"),

        # Legal / Security / Communication
        (["legal", "contract", "agreement", "policy", "terms", "law"], "⚖️"),
        (["security", "protection", "privacy", "safe"], "🛡️"),
        (["password", "locked", "login", "authentication"], "🔐"),
        (["key", "access", "permission"], "🔑"),
        (["rule", "notice", "official"], "📜"),
        (["resume", "cv", "hiring", "recruitment", "candidate", "staff"], "🤝"),
        (["marketing", "campaign", "brand", "promotion", "announcement"], "📣"),

        # Cute / general soft categories
        (["flower", "cute", "aesthetic", "glam", "beautiful"], "🌸"),
        (["nature", "plant", "green"], "🌿"),
        (["star", "special"], "⭐"),
        (["cloud", "soft"], "☁️"),
        (["night", "moon"], "🌙"),
    ]

    for keywords, emoji in checks:
        if any(keyword in name_l for keyword in keywords):
            return emoji

    return "📁"


# Main analyzer: combine extraction, AI/fallback suggestions, and metadata.

def analyze_file(path: str, filename: str, folders: List[Dict] = None) -> Dict:
    """
    Full pipeline. Returns a dict with file_type, extension, word_count,
    keywords, ranked suggestions, top_folder, no_match, and text_preview.
    """
    if folders is None:
        folders = []

    existing_folder_names = [
        f.get("folder") or f.get("name", "")
        for f in folders
        if f.get("folder") or f.get("name")
    ]

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    file_type, raw_text = detect_and_extract(path, filename)

    if file_type == "non-text":
        # Non-text files have no extracted body, so Gemini/TextBlob use the filename.
        raw_text = clean_filename_text(filename)

    if file_type == "empty" or not raw_text.strip():
        raw_text = clean_filename_text(filename)

    limited_text = limit_text(raw_text, 220)
    analysis = ask_gemini_suggestions(filename, limited_text, existing_folder_names)
    ranked = analysis.get("ranked", [])
    summary = _clean_file_summary(analysis.get("summary"), filename, limited_text)

    ai_status = "textblob" if any(s.get("ai_source") == "textblob" for s in ranked) else "gemini"

    return {
        "file_type":    file_type,
        "extension":    ext,
        "word_count":   len(limited_text.split()),
        "keywords":     extract_keywords(limited_text),
        "ranked":       ranked,
        "top_folder":   ranked[0],
        "no_match":     False,
        "text_preview": limited_text[:300],
        "summary":      summary,
        "ai_status":    ai_status,
    }

def analyze_files(file_list: List[Dict], folders: List[Dict] = None) -> List[Dict]:
    """Analyze multiple files. Each item: {"path": "...", "filename": "..."}."""
    results = []
    for item in file_list:
        result = analyze_file(item["path"], item["filename"], folders or [])
        result["filename"] = item["filename"]
        results.append(result)
    return results
